"""Document extraction and reinsertion utilities."""

from __future__ import annotations

import html
import pathlib
import re
from abc import ABC, abstractmethod
from typing import Dict, Iterable, List, Sequence, Tuple

from .errors import UnsupportedFileTypeError, WormholeError
from .structures import TextUnit


RUN_TAG_PATTERN = re.compile(
    r"<run\s+id=\"(?P<id>[^\"]+)\">(?P<content>.*?)</run>",
    re.DOTALL,
)


def _parse_tagged_translation(
    translated: str,
    expected_ids: Sequence[str],
) -> Dict[str, str]:
    """Parse `<run id="">` tagged content and map ids to translated text."""

    mapping: Dict[str, str] = {}
    cursor = 0
    for match in RUN_TAG_PATTERN.finditer(translated):
        prefix = translated[cursor:match.start()]
        if prefix.strip():
            raise WormholeError(
                "Translated output contained unexpected content outside <run> tags."
            )
        run_id = match.group("id")
        if run_id not in expected_ids:
            raise WormholeError(
                f"Translated output contained an unknown run id '{run_id}'."
            )
        if run_id in mapping:
            raise WormholeError(
                f"Translated output duplicated run id '{run_id}'."
            )
        content = html.unescape(match.group("content"))
        mapping[run_id] = content
        cursor = match.end()

    suffix = translated[cursor:]
    if suffix.strip():
        raise WormholeError(
            "Translated output contained unexpected trailing content outside <run> tags."
        )

    if len(mapping) != len(expected_ids):
        missing = [run_id for run_id in expected_ids if run_id not in mapping]
        raise WormholeError(
            "Translation output missing expected runs: "
            + ", ".join(missing)
        )

    return mapping


def _build_units_from_runs(
    runs,
    *,
    unit_prefix: str,
    location: str,
) -> List[TextUnit]:
    """Aggregate paragraph runs into a single translation unit."""

    fragments: List[Tuple[str, object, str]] = []
    for r_idx, run in enumerate(runs):
        text = getattr(run, "text", "")
        if text is None:
            continue
        if not text or not text.strip():
            continue
        fragment_id = f"{unit_prefix}.r{r_idx}"
        fragments.append((fragment_id, run, text))

    if not fragments:
        return []

    if len(fragments) == 1:
        fragment_id, run, text = fragments[0]

        def _single_setter(translated: str) -> None:
            setattr(run, "text", translated)

        return [
            TextUnit(
                unit_id=fragment_id,
                original_text=text,
                setter=_single_setter,
                location=location,
                atomic=False,
            )
        ]

    parts = [
        f'<run id="{fragment_id}">{html.escape(text, quote=False)}</run>'
        for fragment_id, _, text in fragments
    ]
    original_text = "".join(parts)

    fragment_ids = [fragment_id for fragment_id, _, _ in fragments]
    setters = {fragment_id: run for fragment_id, run, _ in fragments}

    def _setter(translated: str) -> None:
        mapping = _parse_tagged_translation(translated, fragment_ids)
        for fragment_id in fragment_ids:
            run = setters[fragment_id]
            new_text = mapping[fragment_id]
            setattr(run, "text", new_text)

    unit = TextUnit(
        unit_id=unit_prefix,
        original_text=original_text,
        setter=_setter,
        location=location,
        atomic=True,
    )
    return [unit]


def _import_docx():
    try:
        from docx import Document  # type: ignore
    except ImportError as exc:  # pragma: no cover - import guard
        raise WormholeError(
            "python-docx is required to process .docx files. "
            "Install the optional dependency with `pip install python-docx`."
        ) from exc
    return Document


def _import_pptx():
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as exc:  # pragma: no cover - import guard
        raise WormholeError(
            "python-pptx is required to process .pptx files. "
            "Install the optional dependency with `pip install python-pptx`."
        ) from exc
    return Presentation


class BaseDocumentHandler(ABC):
    """Common base class for document handlers."""

    def __init__(self, source_path: pathlib.Path):
        self.source_path = source_path
        self.units: List[TextUnit] = []

    @abstractmethod
    def extract_text_units(self) -> List[TextUnit]:
        """Extract translation-ready text units."""

    @abstractmethod
    def save(self, destination: pathlib.Path) -> None:
        """Persist the translated document."""

    def register_units(self, units: Iterable[TextUnit]) -> List[TextUnit]:
        """Store and return the provided units."""

        self.units = list(units)
        return self.units


class DocxDocumentHandler(BaseDocumentHandler):
    """Extracts and reinserts text for Word documents."""

    def __init__(self, source_path: pathlib.Path):
        super().__init__(source_path)
        Document = _import_docx()
        self.document = Document(str(source_path))

    def extract_text_units(self) -> List[TextUnit]:
        units: List[TextUnit] = []
        units.extend(self._extract_body())
        units.extend(self._extract_tables())
        units.extend(self._extract_headers_and_footers())
        return self.register_units(units)

    def save(self, destination: pathlib.Path) -> None:
        self.document.save(str(destination))

    # --- Internal helpers -------------------------------------------------

    def _extract_body(self) -> List[TextUnit]:
        units: List[TextUnit] = []
        for p_idx, paragraph in enumerate(self.document.paragraphs):
            units.extend(
                self._extract_runs(
                    paragraph.runs,
                    unit_prefix=f"body.p{p_idx}",
                    location=f"Body paragraph {p_idx + 1}",
                )
            )
        return units

    def _extract_tables(self) -> List[TextUnit]:
        units: List[TextUnit] = []
        processed_cells = set()
        for t_idx, table in enumerate(self.document.tables):
            for r_idx, row in enumerate(table.rows):
                for c_idx, cell in enumerate(row.cells):
                    cell_key = id(cell._tc)  # type: ignore[attr-defined]
                    if cell_key in processed_cells:
                        continue
                    processed_cells.add(cell_key)
                    base_prefix = (
                        f"body.table{t_idx}.row{r_idx}.cell{c_idx}"
                    )
                    location = (
                        f"Table {t_idx + 1}, row {r_idx + 1}, column {c_idx + 1}"
                    )
                    for p_idx, paragraph in enumerate(cell.paragraphs):
                        units.extend(
                            self._extract_runs(
                                paragraph.runs,
                                unit_prefix=f"{base_prefix}.p{p_idx}",
                                location=location,
                            )
                        )
        return units

    def _extract_headers_and_footers(self) -> List[TextUnit]:
        units: List[TextUnit] = []
        for s_idx, section in enumerate(self.document.sections):
            header = section.header
            footer = section.footer
            for name, container in (("header", header), ("footer", footer)):
                for p_idx, paragraph in enumerate(container.paragraphs):
                    units.extend(
                        self._extract_runs(
                            paragraph.runs,
                            unit_prefix=f"section{s_idx}.{name}.p{p_idx}",
                            location=f"Section {s_idx + 1} {name}",
                        )
                    )
                processed_cells = set()
                for t_idx, table in enumerate(container.tables):
                    for r_idx, row in enumerate(table.rows):
                        for c_idx, cell in enumerate(row.cells):
                            cell_key = id(cell._tc)  # type: ignore[attr-defined]
                            if cell_key in processed_cells:
                                continue
                            processed_cells.add(cell_key)
                            base_prefix = (
                                f"section{s_idx}.{name}.table{t_idx}"
                                f".row{r_idx}.cell{c_idx}"
                            )
                            location = (
                                f"Section {s_idx + 1} {name} table {t_idx + 1}, "
                                f"row {r_idx + 1}, column {c_idx + 1}"
                            )
                            for p_idx, paragraph in enumerate(cell.paragraphs):
                                units.extend(
                                    self._extract_runs(
                                        paragraph.runs,
                                        unit_prefix=f"{base_prefix}.p{p_idx}",
                                        location=location,
                                    )
                                )
        return units

    def _extract_runs(
        self,
        runs,
        *,
        unit_prefix: str,
        location: str,
    ) -> List[TextUnit]:
        return _build_units_from_runs(
            runs,
            unit_prefix=unit_prefix,
            location=location,
        )


class PptxDocumentHandler(BaseDocumentHandler):
    """Extracts and reinserts text for PowerPoint presentations."""

    def __init__(self, source_path: pathlib.Path):
        super().__init__(source_path)
        Presentation = _import_pptx()
        self.presentation = Presentation(str(source_path))

    def extract_text_units(self) -> List[TextUnit]:
        units: List[TextUnit] = []
        units.extend(self._extract_slide_content())
        units.extend(self._extract_notes())
        return self.register_units(units)

    def save(self, destination: pathlib.Path) -> None:
        self.presentation.save(str(destination))

    # --- Internal helpers -------------------------------------------------

    def _extract_slide_content(self) -> List[TextUnit]:
        units: List[TextUnit] = []
        for slide_idx, slide in enumerate(self.presentation.slides):
            slide_prefix = f"slide{slide_idx}"
            for shape_idx, shape in enumerate(slide.shapes):
                base_prefix = f"{slide_prefix}.shape{shape_idx}"
                location = f"Slide {slide_idx + 1}, shape {shape_idx + 1}"
                if getattr(shape, "has_text_frame", False):
                    units.extend(
                        self._extract_text_frame(
                            shape.text_frame,
                            unit_prefix=f"{base_prefix}.tf",
                            location=location,
                        )
                    )
                if getattr(shape, "has_table", False):
                    units.extend(
                        self._extract_table(
                            shape.table,
                            unit_prefix=f"{base_prefix}.table",
                            base_location=location,
                        )
                    )
        return units

    def _extract_notes(self) -> List[TextUnit]:
        units: List[TextUnit] = []
        for slide_idx, slide in enumerate(self.presentation.slides):
            if not getattr(slide, "has_notes_slide", False):
                continue
            notes_slide = slide.notes_slide
            location_base = f"Slide {slide_idx + 1} notes"
            text_frame = notes_slide.notes_text_frame
            if text_frame is not None:
                units.extend(
                    self._extract_text_frame(
                        text_frame,
                        unit_prefix=f"slide{slide_idx}.notes",
                        location=location_base,
                    )
                )
            for shape_idx, shape in enumerate(notes_slide.shapes):
                if getattr(shape, "has_text_frame", False):
                    units.extend(
                        self._extract_text_frame(
                            shape.text_frame,
                            unit_prefix=f"slide{slide_idx}.notes.shape{shape_idx}",
                            location=f"{location_base}, shape {shape_idx + 1}",
                        )
                    )
                if getattr(shape, "has_table", False):
                    units.extend(
                        self._extract_table(
                            shape.table,
                            unit_prefix=f"slide{slide_idx}.notes.shape{shape_idx}.table",
                            base_location=f"{location_base}, table {shape_idx + 1}",
                        )
                    )
        return units

    def _extract_text_frame(
        self,
        text_frame,
        *,
        unit_prefix: str,
        location: str,
    ) -> List[TextUnit]:
        units: List[TextUnit] = []
        for p_idx, paragraph in enumerate(text_frame.paragraphs):
            paragraph_prefix = f"{unit_prefix}.p{p_idx}"
            units.extend(
                _build_units_from_runs(
                    paragraph.runs,
                    unit_prefix=paragraph_prefix,
                    location=location,
                )
            )
        return units

    def _extract_table(
        self,
        table,
        *,
        unit_prefix: str,
        base_location: str,
    ) -> List[TextUnit]:
        units: List[TextUnit] = []
        for r_idx, row in enumerate(table.rows):
            for c_idx, cell in enumerate(row.cells):
                location = (
                    f"{base_location}, row {r_idx + 1}, column {c_idx + 1}"
                )
                prefix = f"{unit_prefix}.row{r_idx}.cell{c_idx}"
                if cell.text_frame is None:
                    continue
                units.extend(
                    self._extract_text_frame(
                        cell.text_frame,
                        unit_prefix=prefix,
                        location=location,
                    )
                )
        return units


def detect_handler(path: pathlib.Path) -> Tuple[str, BaseDocumentHandler]:
    """Select an appropriate handler for the provided file."""

    suffix = path.suffix.lower()
    if suffix == ".docx":
        handler: BaseDocumentHandler = DocxDocumentHandler(path)
        return "docx", handler
    if suffix == ".pptx":
        handler = PptxDocumentHandler(path)
        return "pptx", handler
    raise UnsupportedFileTypeError(
        "This file type isn’t supported — please use .docx or .pptx."
    )
